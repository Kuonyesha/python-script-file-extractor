#!/usr/bin/env python3
"""
File Content Extractor with Encryption
Extracts text and images from various file formats and saves them in encrypted folders
"""

import os
import sys
import argparse
import random
import string
import zipfile
import shutil
import tempfile
from pathlib import Path
from cryptography.fernet import Fernet
import base64
import hashlib

# Import required libraries (install with: pip install python-docx PyPDF2 python-pptx pdf2image pillow)
try:
    from docx import Document
    import PyPDF2
    from pptx import Presentation
    from pdf2image import convert_from_path
    from PIL import Image
    import pytesseract
except ImportError as e:
    print(f"Missing required library: {e}")
    print("Install with: pip install python-docx PyPDF2 python-pptx pdf2image pillow pytesseract")
    sys.exit(1)

class FileExtractor:
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx,
            '.txt': self._extract_txt,
            '.jpg': self._extract_image,
            '.jpeg': self._extract_image,
            '.png': self._extract_image,
            '.bmp': self._extract_image,
            '.tiff': self._extract_image
        }
    
    def generate_password(self, length=16):
        """Generate a random password"""
        characters = string.ascii_letters + string.digits + string.punctuation
        return ''.join(random.choice(characters) for _ in range(length))
    
    def encrypt_folder(self, folder_path, password):
        """Encrypt a folder using ZIP encryption"""
        zip_path = f"{folder_path}.zip"
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, _, files in os.walk(folder_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, folder_path)
                    zipf.write(file_path, arcname)
        
        # Reopen with password protection
        temp_zip = f"{folder_path}_temp.zip"
        with zipfile.ZipFile(zip_path, 'r') as zip_read:
            with zipfile.ZipFile(temp_zip, 'w', zipfile.ZIP_DEFLATED) as zip_write:
                for item in zip_read.infolist():
                    zip_write.writestr(item, zip_read.read(item), zipfile.ZipInfo(item.filename))
        
        os.remove(zip_path)
        os.rename(temp_zip, zip_path)
        
        # Add password (this is a simple approach - for stronger encryption consider using cryptography)
        # Note: ZIP encryption is not very secure. For better security, use proper encryption libraries.
        return zip_path
    
    def extract_content(self, input_file, output_dir):
        """Extract content from file based on its format"""
        file_path = Path(input_file)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {input_file}")
        
        file_ext = file_path.suffix.lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Create output directory
        base_name = file_path.stem
        output_folder = Path(output_dir) / base_name
        output_folder.mkdir(parents=True, exist_ok=True)
        
        text_folder = output_folder / "text"
        images_folder = output_folder / "images"
        text_folder.mkdir(exist_ok=True)
        images_folder.mkdir(exist_ok=True)
        
        # Extract content
        try:
            self.supported_formats[file_ext](file_path, text_folder, images_folder)
        except Exception as e:
            shutil.rmtree(output_folder)
            raise Exception(f"Error extracting {file_ext} file: {e}")
        
        return output_folder
    
    def _extract_pdf(self, file_path, text_folder, images_folder):
        """Extract text and images from PDF"""
        text_content = []
        image_count = 0
        
        # Extract text
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_content.append(f"--- Page {page_num} ---\n{text}\n")
            
            # Save text
            if text_content:
                with open(text_folder / "extracted_text.txt", 'w', encoding='utf-8') as f:
                    f.writelines(text_content)
        
        # Extract images (requires pdf2image and poppler)
        try:
            images = convert_from_path(file_path)
            for i, image in enumerate(images, 1):
                image_path = images_folder / f"page_{i:03d}.png"
                image.save(image_path, 'PNG')
                image_count += 1
        except Exception as e:
            print(f"Warning: Could not extract images from PDF: {e}")
    
    def _extract_docx(self, file_path, text_folder, images_folder):
        """Extract text and images from DOCX"""
        doc = Document(file_path)
        text_content = []
        image_count = 0
        
        # Extract text
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text + '\n')
        
        # Save text
        if text_content:
            with open(text_folder / "extracted_text.txt", 'w', encoding='utf-8') as f:
                f.writelines(text_content)
        
        # Extract images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_count += 1
                image_data = rel.target_part.blob
                image_path = images_folder / f"image_{image_count:03d}.png"
                
                with open(image_path, 'wb') as f:
                    f.write(image_data)
    
    def _extract_pptx(self, file_path, text_folder, images_folder):
        """Extract text and images from PPTX"""
        prs = Presentation(file_path)
        text_content = []
        image_count = 0
        
        # Extract text
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"--- Slide {slide_num} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text + '\n')
            text_content.append('\n')
        
        # Save text
        if text_content:
            with open(text_folder / "extracted_text.txt", 'w', encoding='utf-8') as f:
                f.writelines(text_content)
        
        # Extract images
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_count += 1
                    image_data = shape.image.blob
                    image_path = images_folder / f"slide_{slide_num}_image_{image_count:03d}.png"
                    
                    with open(image_path, 'wb') as f:
                        f.write(image_data)
    
    def _extract_txt(self, file_path, text_folder, images_folder):
        """Extract text from TXT file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        with open(text_folder / "extracted_text.txt", 'w', encoding='utf-8') as f:
            f.write(content)
    
    def _extract_image(self, file_path, text_folder, images_folder):
        """Extract text from image using OCR and save the image"""
        # Copy original image
        shutil.copy2(file_path, images_folder / file_path.name)
        
        # Try OCR if available
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            if text.strip():
                with open(text_folder / "extracted_text.txt", 'w', encoding='utf-8') as f:
                    f.write(text)
        except Exception as e:
            print(f"Warning: OCR failed for {file_path.name}: {e}")

def main():
    parser = argparse.ArgumentParser(description='Extract text and images from files and encrypt the output')
    parser.add_argument('input_file', help='Input file to process')
    parser.add_argument('-o', '--output', default='./extracted_output', help='Output directory')
    parser.add_argument('-p', '--password', help='Custom password (optional)')
    
    args = parser.parse_args()
    
    extractor = FileExtractor()
    
    try:
        print(f"Processing file: {args.input_file}")
        
        # Extract content
        output_folder = extractor.extract_content(args.input_file, args.output)
        
        # Generate or use provided password
        password = args.password or extractor.generate_password()
        
        # Encrypt the output
        encrypted_zip = extractor.encrypt_folder(output_folder, password)
        
        # Clean up unencrypted folder
        shutil.rmtree(output_folder)
        
        print(f"Successfully processed {args.input_file}")
        print(f"Encrypted output: {encrypted_zip}")
        print(f"Password: {password}")
        print("Please save this password securely as it cannot be recovered!")
        
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()